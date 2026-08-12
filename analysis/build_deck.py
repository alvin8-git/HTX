"""Build the HTX biosurveillance briefing deck on the MGI 2025 template.

Data: Basic.stat.xlsx (QC), analysis/species_all.tsv (taxonomy), analysis/amr.tsv (MEGARes),
analysis/vf.tsv (VFDB). Narrative is the conclusion set from docs/biothreat_assessment.md.

SCOPE RULE: every claim on these slides must be derivable from the PFI HTML report alone.
The deck is the briefing form of an engine whose only input is <sample>_en.html, so evidence
that needs the raw FASTQs or an assembly cannot appear here - a reader given the same report
must be able to reach the same conclusion. Removed under this rule (2026-08): unique-read
fraction and the amplification-artifact readings built on it, observed-vs-expected genome GC,
the megahit/abricate assembly cross-check, FASTQ delivery verification, and the unclassified-bin
k-mer probe. Every verdict survived the removal on report-table evidence alone; where a removed
line carried a caveat, the caveat is restated on its HTML-derivable grounds rather than dropped.

Typography/geometry follow HTX_biosurveillance_briefing_modifed.pptx (the user's hand-adjusted
reference): Arial throughout, MINIMUM body size 12 pt - see MIN_PT, which floors every run.

    python3 analysis/build_deck.py    ->  HTX_biosurveillance_briefing.pptx
"""
import csv, json, os, re, sys, openpyxl
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import triage
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
os.chdir(ROOT)
TPL = 'MGI PPT templates_2025_v1.pptx'
OUT = 'HTX_biosurveillance_briefing.pptx'
SAMPLES = ['WBM156', 'WBM174', 'WBM179', 'WBM185', 'WBM232']
MIN_PT = 12.0            # user requirement: nothing smaller than Arial 12
FONT = 'Arial'

BLUE = RGBColor(0x00, 0x53, 0x9B)
LIGHT = RGBColor(0xED, 0xF2, 0xF8)
INK = RGBColor(0x33, 0x33, 0x33)
RED = RGBColor(0xC0, 0x2A, 0x2A)
AMBER = RGBColor(0xB8, 0x6A, 0x00)
GREEN = RGBColor(0x1E, 0x7A, 0x3C)
GREY = RGBColor(0x77, 0x77, 0x77)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SITES = {
    'WBM156': 'Ferry terminal - arrival restroom, tap',
    'WBM174': 'Changi T3 - arrival, automated passport scanner',
    'WBM179': 'Changi T3 - departure, fingerprint reader',
    'WBM185': 'Changi T3 - departure, check-in kiosk touchscreen',
    'WBM232': 'Changi T4 - departure, trolley handles (rows 5-6)',
}

# The flaggable-species list is DERIVED, not hand-written. Two rules, applied identically to all
# five samples:
#   (a) only organisms on the CDC threat list or the WHO/ESKAPE clinical watchlist appear. A
#       near-neighbour used to REFUTE a threat (B. cereus for anthrax, N. subflava for meningitis)
#       is exclusion evidence, not a finding, and belongs on the Category A/B/C slides.
#   (b) resistance genes are never rows of their own. A gene is not an organism; it is evidence
#       about one. Each is attributed as far as the report allows and no further.
LISTED_CAP = 5          # rows before "+N more"; the full set is in the triage HTML report
GENE_CAP = 3
SHORT = {'WHO critical': 'WHO crit', 'WHO high': 'WHO high', 'WHO medium': 'WHO med'}


def badge_for(t):
    """The list badge from the ENGINE's row, not a name lookup. Two watchlist organisms reach the
    list by taxid under a name the rule file does not carry - 'Mycobacteroides abscessus' and
    '[Candida] haemuloni' - so a by-name badge lookup rendered them as '-', i.e. as if they were
    not WHO-listed at all. That is the exact genus-rename drift taxid keying exists to kill,
    reappearing one layer up in the deck."""
    if t['tier'] in ('A', 'B', 'C'):
        return f"CDC {t['tier']}"
    p = (t.get('watch_priority') or '').replace(' (fungal)', '')
    return SHORT.get(p, p) or 'WHO listed'


def escalating(t, attrib):
    """The gene that raises this organism's importance, with the number that qualifies it: what
    share of that gene's candidate-host reads this taxon actually holds. A gene named beside an
    organism reads as that organism's gene; the share is what stops it."""
    hits = []
    for x in attrib:
        c = dict(x['cands']).get(t['taxon'])
        if c is None or not x['pool']:
            continue
        hits.append((x, c / x['pool']))
    if not hits:
        return '-'
    hits.sort(key=lambda h: (0 if h[0]['verdict'] == 'CONFIRM' else 1, -h[0]['breadth']))
    x, sh = hits[0]
    tops = sh >= max(c / x['pool'] for _, c in x['cands'])
    # Never round a real share to 0% - "0% of host pool" reads as "not present", which is the
    # opposite of what a small share means.
    pc = f'{sh:.0%}' if sh >= 0.005 else '<1%'
    return f"{x['group']}: TOPS host pool ({pc})" if tops else f"{x['group']}: {pc} of host pool"


def gene_rank(sm):
    """MEG_ accession -> row number in the PFI drug-resistance table, so a gene on this slide can
    be found by scrolling the source document, exactly as the organism `#` allows."""
    g = triage.load_report(sm)
    return {r['Gene']: i + 1 for i, r in enumerate(g['drugResistance']['DNA']['data'])}


def verdict_driver(sm, listed, attrib):
    """Why this sample carries the verdict it does, in one clause. MONITOR with nothing enriched
    and no acquired gene co-located means something quite different from MONITOR with both, and
    the banner word alone cannot tell them apart."""
    conf = [t for t in listed if t['verdict'] in ('ESCALATE', 'CONFIRM')]
    hc = [x for x in attrib if x.get('high_consequence') and x['verdict'] == 'CONFIRM']
    if conf:
        t = conf[0]
        return (f"{t['taxon']} ({badge_for(t)}) reaches {t['verdict']} - "
                + (f"{len(hc)} high-consequence acquired gene(s) co-located: "
                   f"{', '.join(x['group'] for x in hc)}" if hc
                   else 'acquired resistance of a listed class co-located in this sample'))
    if hc:
        return (f"no listed organism reaches CONFIRM. Driven by {len(hc)} high-consequence gene(s) "
                f"({', '.join(x['group'] for x in hc)}), all topped by commensals")
    enr = [t for t in listed if (t.get('fold') or 0) >= triage.TH['enrichment_fold']]
    if enr:
        return (f"no listed organism reaches CONFIRM; {len(enr)} is site-enriched "
                f"({abbr(enr[0]['taxon'])}) - watch, do not act")
    return (f"no listed organism is site-enriched, and no gene above threshold is topped by one "
            f"({len(listed)} present at background level)")


def abbr(n):
    """Genus to initial: 'Staphylococcus aureus' -> 'S. aureus'. The attribution column is the one
    place where losing the last six characters loses the whole point, so buy them here."""
    parts = n.split()
    # Only binomials. 'Variola virus' -> 'V. virus' is not an abbreviation, it is a different
    # organism's name; the second token has to be a species epithet for the initial to mean anything.
    if len(parts) < 2 or not parts[0][:1].isalpha():
        return n
    if parts[1].lower() in ('virus', 'phage', 'complex', 'group', 'sp.', 'spp.'):
        return n
    return f'{parts[0][0]}. {parts[1]}'


def clip(t, n):
    """Truncate on a word boundary. Mid-word truncation reads as a rendering bug and makes a
    reader distrust the numbers next to it."""
    if len(t) <= n:
        return t
    cut = t[:n].rsplit(' ', 1)[0]
    return cut.rstrip(' ,;-') + ' ...'


def sample_rows(sm):
    """(listed organisms, gene attributions) for one sample, straight from the triage engine."""
    g = triage.load_report(sm)
    _, clas = triage.gate_integrity(sm, g)
    reps = {x: triage.load_report(x) for x in SAMPLES}
    cl = {}
    for x, gg in reps.items():
        _, cl[x] = triage.gate_integrity(x, gg)
    loads = triage.loads_by_taxon(reps, cl)
    genes = triage.triage_genes(g)
    taxa = triage.triage_taxa(sm, g, loads, genes, True)
    sp = g['indentification_DNA']['speciesData']['data']
    rank = {r['Scientific Name']: i + 1 for i, r in enumerate(sp)}
    present = {r['Scientific Name']: int(r['Real Read']) for r in sp}

    order = {'ESCALATE': 0, 'CONFIRM': 1, 'MONITOR': 2, 'NO_ACTION': 3}
    # NOT_TESTED rows with reads are kept: since 2026-08-12 that tier also covers "confirmatory
    # marker not assessable at this coverage", which applies to organisms that ARE present and
    # abundant (S. aureus at 3,654 reads in WBM232). The zero-read RNA agents are excluded by
    # t['real'] alone, which is what the old NOT_TESTED filter was really for.
    listed = [t for t in taxa if t['tier'] not in ('-', '') and t['real']]
    # Ordered by ABUNDANCE, which is how a reader scans a species list - not by verdict, which
    # put a 473-read organism above a 2,300-read one and read as a claim about which matters more.
    # Any organism the engine escalated is pulled in regardless of where abundance puts it.
    listed.sort(key=lambda t: -t['real'])
    top = listed[:LISTED_CAP]
    hot = [t for t in listed if t['verdict'] in ('ESCALATE', 'CONFIRM') and t not in top]
    shown = (hot + top)[:LISTED_CAP] if hot else top
    shown.sort(key=lambda t: -t['real'])
    return listed, gene_attribution(genes, present), rank, shown


def short_why(t):
    """One line, about THIS organism only, built from its own fields rather than by slicing the
    engine's prose - a clause cut out of a sentence about the sample reads as a claim about the
    row it lands next to, which is the error this whole table exists to remove."""
    f, w = t.get('fold'), t['why']
    if f == float('inf'):
        base = 'Detected only in this sample'
    elif f and f >= triage.TH['enrichment_fold']:
        base = f'{f:.1f}x enriched vs other sites'
    elif f:
        base = f'Not site-enriched ({f:.2f}x)'
    else:
        base = 'Above the read floor'
    if 'below min' in w:
        base = 'Below the 50-read floor'
    if 'ABSENT - downgraded' in w:
        base += '; marker absent'
    elif 'marker(s)' in w and 'PRESENT' in w:
        base += '; MARKER PRESENT'
    if 'estimate inflated' in w:
        base += '; est. inflated'
    if 'kitome' in w:
        base += '; reagent genus'
    return base


def gene_attribution(genes, present):
    """How far the report allows a resistance gene to be attributed to an organism. Four states,
    and the distinction that matters is intrinsic vs acquired:

      species      exactly one documented host of this gene is in the sample
      genus        several, all one genus - the genus is settled, the species is not
      candidates   several genera - name the most abundant and say how many compete
      unattributed no documented host present at all: a gene with no organism to hang it on

    An INTRINSIC gene is the strong case. It is part of its genus's own chromosome, so if that
    genus is present the gene came from it - the only open question is which species. An ACQUIRED
    gene is the weak case: it is mobile by definition, which is exactly why it cannot be placed.
    """
    hints = {k: v for k, v in triage.RULES.get('amr_host_hints', {}).items() if k != '_comment'}
    out = []
    for x in genes:
        if x['verdict'] == 'NO_ACTION':
            continue
        h = hints.get(x['group'])
        intrinsic = x['class'] in ('intrinsic', 'core_essential', 'regulator')
        cands = []
        if h:
            cands = sorted(((n, c) for n, c in present.items()
                            if any(n.startswith(t + ' ') or n == t for t in h['taxa']) and c),
                           key=lambda kv: -kv[1])
        genera = {n.split()[0] for n, _ in cands}
        tot = sum(c for _, c in cands)
        # The decision-relevant number is not "who is top" but "what share does the organism a
        # reader is about to name actually hold". Naming S. aureus for mecA when it holds 0.7% of
        # the staphylococcal reads is the specific error this column exists to stop.
        lst = next(((n, c) for n, c in cands if list_badge(n)), None)
        badge = SHORT.get(list_badge(lst[0]), list_badge(lst[0])) if lst else ''
        if not cands:
            state = 'unattributed'
            who = 'No documented host of this gene is present in this sample'
        elif len(cands) == 1:
            state = 'species'
            who = f'{abbr(cands[0][0])} - the only documented host present'
        else:
            state = 'genus' if len(genera) == 1 else 'candidates'
            top, tc = cands[0]
            pc = lambda v: f'{v:.0%}' if v >= 0.005 else '<1%'
            who = f'{abbr(top)} holds {pc(tc/tot)} of the host pool ({len(cands)} spp.)'
            if lst and lst[0] == top:
                who = f'{abbr(top)} ({badge}) tops the host pool - {pc(tc/tot)} of {len(cands)} spp.'
            elif lst:
                who += f'; listed {abbr(lst[0])} only {pc(lst[1]/tot)}'
        hits_listed = bool(lst and cands and lst[0] == cands[0][0])
        out.append({**x, 'state': state, 'who': who, 'intrinsic': intrinsic,
                    'hits_listed': hits_listed, 'cands': cands, 'pool': tot,
                    'n_cands': len(cands)})
    rank = {'species': 0, 'genus': 1, 'candidates': 2, 'unattributed': 3}
    # high_consequence FIRST, before breadth. Sorting a briefing table by coverage puts a 98%
    # commensal erm gene above mecA at 90.9% - technically ranked, operationally backwards.
    # A gene that CAN be pinned to a listed organism outranks one that cannot, whatever its
    # coverage. An attributable intrinsic gene is the most actionable row on the slide, and it was
    # being sorted off the table by commensal genes with higher breadth.
    # Verdict outranks everything: high_consequence marks which gene matters IF real, not whether
    # it is real, so letting it jump a tier puts a MONITOR-grade call above three CONFIRMs.
    out.sort(key=lambda r: (0 if r['verdict'] == 'CONFIRM' else 1,
                            0 if r.get('high_consequence') else 1,
                            0 if r['hits_listed'] else 1,
                            rank[r['state']], -r['breadth']))
    return out


# Every string here renders at Arial 12 minimum (MIN_PT floors _style), so the box holds about
# five lines. These are written to that budget rather than trimmed to it afterwards.
ACTIONABLE = {
    'WBM156': ('MONITOR', AMBER,
               'Human oral/salivary flora plus water-associated organisms - what a restroom tap should '
               'look like. No threat agent; 13 resistance classes, none acquired and of consequence.\n'
               'WHY "MONITOR" NOT "NO ACTION": P. rettgeri (WHO critical) is seen in NO other swab - '
               '310 rpm against zero elsewhere - so it is site-specific, not batch background. It '
               'carries no gene of consequence, hence watch-and-repeat rather than act.'),
    'WBM174': ('NO ACTION', GREEN,
               'Skin-flora dominated (C. acnes 34.2%) - a hand-touched surface behaving normally. The '
               'staphylococcal genes above are community-normal carriage of the genus, not of S. aureus, '
               'which holds 0.7% of the staphylococcal reads.\n'
               'WHY "NO ACTION": no listed organism is site-enriched, and every gene above threshold is '
               'topped by a commensal. Botulism is refuted on organism (11 reads) and toxin (bont absent).'),
    'WBM179': ('NO ACTION', GREEN,
               'Deepest classified dataset (3.1M reads) and nothing threat-level. Oral + skin flora from '
               'fingers; 21 AMR classes, no mecA and no ESBL. mecI without mecA is a regulator, not '
               'resistance.\n'
               'WHY "NO ACTION": K. aerogenes at 3.0x is below the 5x enrichment bar, and every gene '
               'above threshold is topped by a commensal. Best baseline in the batch.'),
    'WBM185': ('INVESTIGATE', RED,
               'Highest AMR burden in the batch: 90 genes / 21 classes, mecA at 90.9% and CTX-M at 82.9%.\n'
               'ACTION: resolve the mecA host by culture with AST (antimicrobial susceptibility testing). '
               'It is staphylococcal, but S. aureus holds 0.7% of the staphylococcal reads against '
               'S. hominis at 52% - calling this MRSA assigns the gene to one of the rarest candidates. '
               'No listed ORGANISM drives this verdict; the three genes do.'),
    'WBM232': ('INVESTIGATE', RED,
               'The one operationally significant finding. A. baumannii at 4.72%, load 6-43x above every '
               'other site - real enrichment, not reagent background. It TOPS the host pool of both '
               'acquired genes (21% of 63 candidates): best-supported host, not a proven one.\n'
               'ACTION: re-swab T4 trolley handles WITH A CULTURE ARM - metagenomics cannot give an '
               'antibiogram, and lpxA presence is not colistin resistance without a mutation call. Not '
               'ESCALATE: that needs a CDC agent with its marker; A. baumannii is WHO-critical.'),
}

CAT_A = [
    ('Bacillus anthracis', 'NEGATIVE',
     'B. cereus group present (WBM174 22 rds, WBM185 43 rds). '
     'Excluded at PLASMID level: zero pXO1 (pagA/lef/cya/atxA), zero pXO2 (capA-E). '
     'Only VFDB hits are chromosomal isdC and GBAA_RS23245, shared across the whole B. cereus group.'),
    ('Clostridium botulinum', 'NEGATIVE',
     'WBM174: 11 species-specific reads, far below the 50-read floor. '
     'No botulinum neurotoxin (bont) gene in any sample.'),
    ('Yersinia pestis', 'NEGATIVE',
     'No Yersinia of any species in any sample. The ybtT/irp2 hits are yersiniabactin siderophore genes '
     'carried by K. pneumoniae and E. coli - a shared iron-uptake island, not plague.'),
    ('Francisella tularensis', 'NEGATIVE', 'Absent, and no Francisella congeners detected.'),
    ('Variola / Orthopoxvirus', 'NEGATIVE',
     'Meaningful negative - these are dsDNA viruses and WOULD have been sequenced by this DNA library. '
     'The call rests on taxonomy against vaccinia / cowpox / mpox: VFDB is a bacterial database, so '
     'there is no orthopoxvirus marker gene to confirm with.'),
    ('Viral haemorrhagic fevers\n(Ebola, Marburg, Lassa, Junin)', 'NOT ASSESSABLE',
     'All are RNA viruses. There is NO RNA library in this batch, so they are structurally undetectable. '
     'This is a missing test, not a negative result.'),
]

CAT_BC = [
    ('B', 'Brucella spp.', '"Brucella anthropi" in all 5 samples (14-1,824 rds)',
     'NOT BRUCELLOSIS - this is Ochrobactrum anthropi, renamed into Brucella in 2020. Ubiquitous reagent '
     'contaminant, carrying a false "Human Infection: Y" flag in every sample. Now separable by gene, '
     'not just by name: btpA/btpB are Brucella TIR effectors that Ochrobactrum lacks. Absent here.'),
    ('B', 'Vibrio cholerae', 'WBM179, 11 reads', 'REFUTED - below the 50-read floor, no cholera toxin genes.'),
    ('B', 'Clostridium perfringens', 'WBM174 84, WBM185 248, WBM232 29 rds',
     'Present at trace level. Ubiquitous soil/gut anaerobe; no epsilon-toxin (etx) gene detected.'),
    ('B', 'Escherichia coli', '62-118 real reads (est. up to 12,076)',
     'Trace, and the Bracken estimate is inflated ~100x. No O157:H7 markers (stx1/stx2/eae) detected.'),
    ('B', 'Salmonella / Shigella / Coxiella /\nBurkholderia mallei / pseudomallei / C. psittaci', 'ABSENT',
     'No reads assigned to any of these in any sample.'),
    ('B', 'Rickettsia felis / massiliae', 'WBM174 101, WBM185 151/53 rds',
     'Low-confidence trace calls; R. prowazekii (the listed agent) is absent.'),
    ('C', 'Mycobacterium tuberculosis (MDR-TB)', 'ABSENT',
     'No M. tuberculosis complex reads. The Mycobacterium present (M. avium, M. grossiae, M. paragordonae, '
     'M. intracellulare) are environmental NTM, typical of building water systems. Backed by esxA/esxB '
     '(RD1) - deleted in BCG and absent from most NTM. Absent here.'),
    ('C', 'Nipah, Hantavirus, TBE, Yellow fever', 'NOT ASSESSABLE',
     'All RNA viruses - no RNA library. Same structural blind spot as the Category A VHFs.'),
]


# --- helpers ---------------------------------------------------------------
def _style(run, size, bold, color):
    run.font.size = Pt(max(size, MIN_PT))
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT


def txbox(slide, l, t, w, h, text, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT,
          space_after=4, line=1.0):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line
        r = p.add_run(); r.text = ln
        _style(r, size, bold, color)
    return tb


def title(slide, text, sub=None, w=12.4, size=26):
    txbox(slide, 0.5, 0.28, w, 0.55, text, size=size, bold=True, color=BLUE)
    if sub:
        txbox(slide, 0.5, 0.88, 12.4, 0.3, sub, size=12, color=GREY)
    ln = slide.shapes.add_shape(1, Inches(0.5), Inches(1.22), Inches(12.33), Inches(0.028))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE
    ln.line.fill.background(); ln.shadow.inherit = False


def table(slide, l, t, w, rows, widths, sizes=(12, 12), header_h=0.34, row_h=0.30,
          colors=None, bolds=None):
    nr, nc = len(rows), len(rows[0])
    h = header_h + row_h * (nr - 1)
    gt = slide.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(h)).table
    gt.rows[0].height = Inches(header_h)
    for i in range(1, nr):
        gt.rows[i].height = Inches(row_h)
    tot = sum(widths)
    for c, ww in enumerate(widths):
        gt.columns[c].width = Inches(w * ww / tot)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = gt.cell(r, c)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if r == 0 else (LIGHT if r % 2 == 0 else WHITE)
            numeric = str(val).replace('.', '').replace('%', '').replace(',', '').replace('x', '')
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.RIGHT if (c > 0 and r > 0 and numeric.isdigit()) else PP_ALIGN.LEFT
                for run in p.runs:
                    _style(run, sizes[0] if r == 0 else sizes[1],
                           (r == 0) or (bolds or {}).get((r, c), False),
                           WHITE if r == 0 else (colors or {}).get((r, c), INK))
    return gt


def chip(slide, l, t, w, h, text, rgb, size=13):
    s = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb
    s.line.fill.background(); s.shadow.inherit = False
    s.text_frame.word_wrap = True
    p = s.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    _style(r, size, True, WHITE)
    return s


def panel(slide, l, t, h, rgb=BLUE):
    """Thin left rule used to group a block of text."""
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(0.07), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb
    s.line.fill.background(); s.shadow.inherit = False
    return s


# --- data loading ----------------------------------------------------------
def qc():
    return {s: {k: v for k, v in openpyxl.load_workbook(f'{s}/Basic.stat.xlsx').active
                .iter_rows(values_only=True)} for s in SAMPLES}


def species():
    rows = [x for x in csv.DictReader(open('analysis/species_all.tsv'), delimiter='\t')
            if x['level'] == 'speciesData']
    return {s: sorted([x for x in rows if x['sample'] == s], key=lambda x: -float(x['ab']))
            for s in SAMPLES}


def amr_counts():
    a = list(csv.DictReader(open('analysis/amr.tsv'), delimiter='\t'))
    return {s: (len([x for x in a if x['sample'] == s]),
                len({x['Class'] for x in a if x['sample'] == s})) for s in SAMPLES}


def vf_counts():
    v = list(csv.DictReader(open('analysis/vf.tsv'), delimiter='\t'))
    return {s: len([x for x in v if x['sample'] == s]) for s in SAMPLES}


# --- slides ----------------------------------------------------------------
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[3])


def s_cover(prs):
    cov = prs.slides[0]
    for sh in cov.shapes:
        if sh.name == '副标题 2':
            sh.text_frame.text = 'Biosurveillance of Singapore Transport Hubs'
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(30); r.font.bold = True; r.font.name = FONT
        if sh.name == '文本框 6':
            sh.text_frame.text = 'DNA shotgun metagenomics  |  5 surface swabs  |  2026.07.31'
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12); r.font.name = FONT
    txbox(cov, 0.82, 4.55, 7.6, 1.1,
          'Threat-agent screen, CDC Category A/B/C assessment, and AMR profile\n'
          'Kraken2 / Bracken taxonomy - MEGARes AMR - VFDB virulence',
          size=13, color=GREY, line=1.15)


def s_method(prs):
    s = blank(prs)
    title(s, 'The assay, and what it can and cannot see')
    txbox(s, 0.5, 1.45, 6.0, 0.3, 'METHOD', size=13, bold=True, color=BLUE)
    txbox(s, 0.5, 1.8, 6.0, 3.9,
          'DNA-only shotgun metagenomics, 150 bp paired-end\n'
          'Host depletion, then Kraken2 classification with Bracken abundance re-estimation\n'
          'AMR called against MEGARes; virulence factors against VFDB\n'
          'Reports generated by PFI software v5.1.2 / database v5.1.1\n\n'
          'Scope of this briefing:\n'
          'Every claim on these slides is derived from the PFI HTML report alone - the same '
          'document you were sent. Nothing here depends on the raw FASTQs or on a de novo '
          'assembly, so any conclusion drawn can be checked against the report itself.',
          size=14, line=1.25, space_after=6)
    txbox(s, 6.9, 1.45, 6.0, 0.3, 'THREE THINGS THAT WILL MISLEAD YOU', size=13, bold=True, color=RED)
    txbox(s, 6.9, 1.8, 6.0, 3.9,
          '1. Judge on Real Read, not Abundance. Abundance and Estimate Read are Bracken '
          'redistributions and can be wildly inflated - S. agalactiae in WBM174 has 29 species-specific '
          'reads but an estimate of 8,301.\n\n'
          '2. A gene hit is not an organism. MEGARes and VFDB rows carry no organism column, so no '
          'resistance or virulence gene in this report is attributable to any species in it.\n\n'
          '3. The "Human Infection: Y" flag is unreliable at the edges - it fires on renamed '
          'environmental organisms.',
          size=14, line=1.25, space_after=6)
    txbox(s, 0.5, 5.85, 12.4, 1.2,
          'Discriminators applied throughout: (a) Real Read against a 50-read floor, judged separately '
          'from the Bracken estimate; (b) depth-normalised load (reads per million classified) across '
          'all five sites - a real site finding is enriched in one sample, a reagent contaminant is flat '
          'everywhere; (c) breadth of the reference gene rather than depth, so a conserved fragment read '
          'deeply does not pass as a whole gene; (d) for B. anthracis, presence of the pXO1/pXO2 plasmid '
          'markers rather than chromosomal genes shared across the B. cereus group.',
          size=12, color=GREY, line=1.2)


def s_qc(prs, Q):
    s = blank(prs)
    title(s, 'Initial QC - all five samples', 'Basic.stat.xlsx, DNA libraries, 150 bp PE')
    rows = [['Sample', 'Site', 'Raw reads', 'GC', 'Q30', 'Low-qual', 'Host', 'Clean', 'Unclass.', 'Classified']]
    colors, bolds = {}, {}
    for i, sm in enumerate(SAMPLES, start=1):
        q = Q[sm]
        pct = lambda k: q[k].split(' ')[1].strip('()')
        rows.append([sm, SITES[sm].split(' - ')[0], f"{q['Raw_Read']:,}", q['Raw_GC'], q['Raw_Q30'],
                     pct('Lowquality_Read'), pct('Host_Read'), pct('Clean_Read'),
                     pct('Unclassified_Read'), f"{int(q['Classified_Read'].split(' ')[0]):,}"])
        colors[(i, 0)] = BLUE
    colors[(1, 6)] = colors[(1, 9)] = RED
    bolds[(1, 6)] = bolds[(1, 9)] = True
    table(s, 0.5, 1.5, 12.33, rows, [1.1, 2.74, 1.42, 0.82, 0.88, 1.1, 0.99, 0.99, 1.04, 1.26],
          sizes=(12, 14), row_h=0.38, colors=colors, bolds=bolds)
    txbox(s, 0.5, 3.9, 12.4, 0.3, 'READ THIS TABLE FROM RIGHT TO LEFT', size=13, bold=True, color=BLUE)
    txbox(s, 0.5, 4.25, 12.4, 2.55,
          'Base quality is uniformly excellent - Q30 98.0-98.5%, low-quality reads under 1.5% in every '
          'sample. Quality is not the limiting factor here; BIOMASS is.\n\n'
          'The number that determines sensitivity is the last column, classified reads. WBM156 lost '
          '91.3% of its reads to host depletion and ended with 622,642 classified reads - roughly 5x '
          'less microbial data than WBM179. Its shorter hit list is a detection-limit artifact, not '
          'evidence of a cleaner tap.\n\n'
          'Unclassified fractions of 72-90% are normal for environmental swabs against a clinical '
          'database, but they are the batch\'s largest blind spot and the report says nothing about '
          'what is in them. Anything absent from the PFI database lands here and is invisible to every '
          'slide that follows.',
          size=14, line=1.25, space_after=7)


def s_caveats(prs):
    s = blank(prs)
    title(s, 'QC caveats that change how you read the results')
    items = [
        ('No negative control', AMBER, 1.5,
         'No blank extraction control was included in this batch. Trace calls below ~100 reads cannot be '
         'formally separated from reagent background. We compensated with a cross-sample kitome analysis '
         '(216 core taxa shared by all 5 samples; 50 known contaminant genera; 46 with abundance-vs-depth '
         'correlation rho <= -0.6, e.g. P. putida at rho = -1.00) - but a blank would have been better.'),
        ('A read count is not a molecule count', AMBER, 1.5,
         'The report gives read counts only. It cannot show whether a taxon\'s reads are distinct '
         'molecules or repeated copies of one amplified fragment, and both look identical in the table. '
         'Trace calls are therefore refuted on the 50-read floor and on marker absence - never on '
         'molecule counts, which this document does not contain.'),
        ('AMR genes are called off reads', RED, 1.5,
         'The PFI pipeline reports resistance genes without linking them to an organism, and MEGARes '
         'holds many near-identical alleles of the same gene, so ONE true gene lights up SEVERAL MEG_ '
         'accessions. Both effects are unpacked on the WBM185 and WBM232 evidence slides. This is the '
         'single biggest interpretive trap in the standard report.'),
    ]
    y = 1.5
    for head, rgb, h, body in items:
        panel(s, 0.5, y, h, rgb)
        txbox(s, 0.72, y, 2.6, 0.3, head, size=13.5, bold=True, color=rgb)
        txbox(s, 3.3, y - 0.02, 9.5, h, body, size=14, line=1.2)
        y += h + 0.14


def s_topspecies(prs, SP):
    s = blank(prs)
    title(s, 'Top species by abundance - side by side',
          'Bracken abundance %; species-specific ("real") reads in brackets')
    rows = [['#'] + SAMPLES]
    for i in range(6):
        rows.append([str(i + 1)] + [f"{SP[sm][i]['name']}\n{SP[sm][i]['ab']}%  "
                                    f"({int(SP[sm][i]['real']):,} rds)" for sm in SAMPLES])
    table(s, 0.5, 1.5, 12.33, rows, [0.35, 2.4, 2.4, 2.4, 2.4, 2.4], sizes=(12, 12), row_h=0.62)
    txbox(s, 0.5, 5.65, 12.4, 1.1,
          'Cutibacterium acnes - skin sebaceous flora - tops four of five samples and reaches 57.4% in '
          'WBM232. These are hand-touched surfaces, so a human skin signature dominating is the expected '
          'baseline, not a finding. WBM185 is the exception: Staphylococcus and Kocuria displace C. acnes, '
          'and WBM232 is the only sample where a recognised nosocomial pathogen (A. baumannii) reaches '
          'the top three.', size=12, color=GREY, line=1.2)


def list_badge(label):
    """'CDC A' / 'WHO critical' / None for a flagged-row label, read from triage_rules.json.

    Taken from the rule file rather than typed onto the slide so the deck cannot drift from the
    engine. Returns None for gene rows ('mecA - MEG_3778') and for species on neither list -
    absence of a badge is itself information: the row was flagged for evidence, not for
    membership. Handles the abbreviated forms used on crowded rows ('L. pneumophila')."""
    rules = json.load(open(os.path.join(ROOT, 'analysis', 'triage_rules.json')))
    tl = {k: v for k, v in rules['threat_list'].items() if k != '_comment'}
    wl = {k: v for k, v in rules['clinical_watchlist'].items() if k != '_comment'}
    if 'MEG_' in label:
        return None
    name = label.strip()
    if name in tl:
        return 'CDC ' + tl[name]['tier']
    if name in wl:
        return wl[name]['priority']
    m = re.match(r'^([A-Z])\.\s+(\S+)$', name)          # 'L. pneumophila'
    if m:
        for k, v in list(tl.items()) + list(wl.items()):
            parts = k.split(' ')
            if len(parts) >= 2 and parts[0][0] == m.group(1) and parts[1] == m.group(2):
                return 'CDC ' + v['tier'] if k in tl else v['priority']
    return None


def s_sample(prs, sm, Q, SP, AM, VF):
    s = blank(prs)
    verdict, vcolor, action = ACTIONABLE[sm]
    title(s, f'{sm}  -  {SITES[sm]}', w=10.1, size=22)
    chip(s, 10.88, 0.3, 1.95, 0.45, verdict, vcolor)
    q = Q[sm]
    txbox(s, 0.5, 1.4, 12.35, 0.2,
          f"{q['Raw_Read']:,} raw reads   |   host {q['Host_Read'].split(' ')[1].strip('()')}   |   "
          f"{int(q['Classified_Read'].split(' ')[0]):,} classified   |   {len(SP[sm])} species   |   "
          f"{AM[sm][0]} AMR genes / {AM[sm][1]} classes   |   {VF[sm]} VF hits",
          size=12, color=BLUE, bold=True)

    listed, attrib, rank, shown = sample_rows(sm)
    grank = gene_rank(sm)

    # One line naming what actually produced the banner verdict. Without it MONITOR reads as a
    # judgement; with it, MONITOR is a rule the reader can disagree with.
    txbox(s, 0.5, 1.62, 12.33, 0.22, 'VERDICT DRIVER:  ' + verdict_driver(sm, listed, attrib),
          size=11.5, bold=True, color=vcolor)

    txbox(s, 0.5, 1.90, 12.33, 0.28, 'CDC / WHO LISTED ORGANISMS PRESENT', size=12.5, bold=True,
          color=RED)
    vc = {'ESCALATE': RED, 'CONFIRM': RED, 'MONITOR': AMBER, 'NO_ACTION': GREY}
    rows = [['#', 'Organism', 'List', 'Reads (ab.)', 'Assessment', 'Co-located gene (MEGARes)']]
    colors, bolds = {}, {}
    for i, t in enumerate(shown, start=1):
        rows.append([str(rank.get(t['taxon'], '')), t['taxon'], badge_for(t),
                     f"{t['real']:,} ({t['abundance']})", clip(short_why(t), 52),
                     clip(escalating(t, attrib), 48)])
        colors[(i, 2)] = BLUE
        colors[(i, 1)] = vc.get(t['verdict'], INK)
        colors[(i, 5)] = AMBER
        bolds[(i, 1)] = t['verdict'] in ('ESCALATE', 'CONFIRM')
    table(s, 0.5, 2.18, 12.33, rows, [0.42, 2.45, 0.92, 1.32, 3.35, 3.87], sizes=(11, 11),
          row_h=0.32, colors=colors, bolds=bolds)
    n_more = max(0, len(listed) - len(shown))
    txbox(s, 0.5, 4.12, 12.33, 0.2,
          (f'+{n_more} further listed - full set in the triage report.   ' if n_more else '')
          + 'Column meanings on the "How to read a sample slide" page.',
          size=10, color=GREY)

    txbox(s, 0.5, 4.42, 12.33, 0.28, 'RESISTANCE GENES  -  AND WHO THEY CAN BE ATTRIBUTED TO',
          size=12.5, bold=True, color=BLUE)
    rows = [['#', 'Gene', 'Source', 'Cov / Depth', 'Most likely host organism']]
    colors, bolds = {}, {}
    scol = {'species': GREEN, 'genus': AMBER, 'candidates': AMBER, 'unattributed': GREY}
    for i, x in enumerate(attrib[:GENE_CAP], start=1):
        rows.append([str(grank.get(x['allele'], '')), x['group'], 'MEGARes',
                     f"{x['breadth']:.1f}% / {x['depth']:.1f}x", clip(x['who'], 84)])
        colors[(i, 4)] = scol[x['state']]
        colors[(i, 0)] = colors[(i, 1)] = RED if x['verdict'] == 'CONFIRM' else INK
        bolds[(i, 1)] = x['verdict'] == 'CONFIRM'
    table(s, 0.5, 4.70, 12.33, rows, [0.42, 1.35, 1.0, 1.45, 8.11], sizes=(11, 11),
          row_h=0.32, colors=colors, bolds=bolds)
    n_gm = max(0, len(attrib) - GENE_CAP)
    txbox(s, 0.5, 6.02, 12.33, 0.2,
          'BOLD RED = CONFIRM; plain = MONITOR.   All rows MEGARes; VFDB markers appear above.'
          + (f'   +{n_gm} more above threshold.' if n_gm else ''), size=10, color=GREY)

    txbox(s, 0.5, 6.28, 12.33, 0.26, f'WHAT IS ACTIONABLE  -  {verdict}',
          size=12.5, bold=True, color=vcolor)
    txbox(s, 0.5, 6.52, 12.33, 0.95, action, size=10.5, line=1.06, space_after=1)


def s_howtoread(prs):
    s = blank(prs)
    title(s, 'How to read a sample slide',
          'The same columns, the same rules, on all five samples')
    txbox(s, 0.5, 1.36, 12.33, 0.26, 'THREE DIFFERENT PERCENTAGES APPEAR. THEY HAVE DIFFERENT '
          'DENOMINATORS.', size=12.5, bold=True, color=RED)
    rows = [['Number', 'Denominator', 'Example: WBM185, S. maltophilia / CTX-M'],
            ['Abundance %  (organism)', 'Species reads / all classified reads',
             '0.69%  -  2,090 of 2,006,566 classified reads'],
            ['Coverage %  (gene)', 'Bases of the REFERENCE GENE covered / gene length. Says nothing '
             'about how much organism is present', '82.9%  -  of the CTX-M reference sequence'],
            ['Host-pool share %  (link)', "This organism's reads / reads of EVERY documented host of "
             'that gene here. THE number for how likely the gene is this organism\'s',
             '5.9%  -  of the 77 organisms that could carry CTX-M']]
    table(s, 0.5, 1.68, 12.33, rows, [2.4, 5.4, 4.53], sizes=(11.5, 11), row_h=0.64,
          colors={(i, 0): BLUE for i in (1, 2, 3)}, bolds={(i, 0): True for i in (1, 2, 3)})

    txbox(s, 0.5, 4.28, 12.33, 0.26, 'WHAT THE OTHER COLUMNS MEAN', size=12.5, bold=True, color=BLUE)
    rows = [['Column', 'Meaning'],
            ['#', "Row number in the PFI report's own table - scroll to it to check any line"],
            ['Bold red', 'CONFIRM: real and actionable. Plain text is MONITOR'],
            ['Not site-enriched', 'Depth-normalised load is NOT >=5x every other swab - present '
             'batch-wide, so background rather than a site event'],
            ['Co-located gene  /\nMost likely host', "The strongest gene whose documented host range "
             "covers this organism's genus, and the share of that gene's host pool the organism "
             'holds. A HIGH share escalates the organism; a LOW share is why a gene sitting beside a '
             'name is not that organism\'s gene'],
            ['Two databases', 'MEGARes (drug resistance) has NO organism column - hence a host pool '
             'instead of a species. VFDB (virulence) DOES name a reference pathogen, so virulence '
             'markers report against the organism and appear in Assessment, not here']]
    table(s, 0.5, 4.58, 12.33, rows, [2.1, 10.23], sizes=(11.5, 11), row_h=0.46,
          colors={(i, 0): BLUE for i in range(1, 6)},
          bolds={(i, 0): True for i in range(1, 6)})


def s_botulism(prs):
    s = blank(prs)
    title(s, 'WBM174 - botulism, and why the bont gene is the actual test',
          'Why 11 reads of Clostridium botulinum is not a botulism finding')
    txbox(s, 0.5, 1.42, 6.05, 0.28, 'WHAT MAKES BOTULISM A CATEGORY A AGENT', size=13, bold=True, color=BLUE)
    panel(s, 0.5, 1.75, 2.15, BLUE)
    txbox(s, 0.72, 1.75, 5.83, 2.15,
          'The threat is not the bacterium, it is the protein. Botulinum neurotoxin (BoNT) is the most '
          'potent toxin known - lethal dose in the low nanograms per kilogram. It blocks acetylcholine '
          'release at the neuromuscular junction, causing descending flaccid paralysis and respiratory '
          'failure.\n\n'
          'It is a Category A agent because it can be disseminated in food or as an aerosol, and because '
          'even a small cluster consumes ICU beds, ventilators and antitoxin stocks.',
          size=13, line=1.2, space_after=6)
    txbox(s, 6.78, 1.42, 6.05, 0.28, 'WHY THE ORGANISM ALONE MEANS NOTHING', size=13, bold=True, color=BLUE)
    panel(s, 6.78, 1.75, 2.15, BLUE)
    txbox(s, 7.0, 1.75, 5.83, 2.15,
          'C. botulinum is a spore-forming anaerobe living in soil and dust worldwide, so traces on an '
          'environmental surface are unremarkable.\n\n'
          'Toxin production comes from the bont gene cluster (types A-G; A, B, E, F cause human disease), '
          'which is MOBILE - chromosome, plasmid or prophage depending on strain. Non-toxigenic '
          'C. botulinum exists, and toxigenic C. butyricum and C. baratii exist. Species identity '
          'therefore does not imply toxin.',
          size=13, line=1.2, space_after=6)

    txbox(s, 0.5, 4.4, 12.33, 0.28, 'WHAT THE WBM174 DATA ACTUALLY SHOWS', size=13, bold=True, color=GREEN)
    rows = [['Test', 'Result', 'Interpretation'],
            ['Species-specific reads', '11', 'Below any credible detection threshold to begin with'],
            ['Estimate Read', '37', 'Even the inflated Bracken estimate stays far below the floor'],
            ['bont toxin gene (VFDB)', '0 reads', 'Zero in WBM174 and zero in all four other samples'],
            ['Other Clostridium virulence hits', 'C. perfringens only',
             'fbpA, tadA, nagH/nagI adhesins in WBM185 - no neurotoxin of any kind anywhere in the batch']]
    colors = {(i, 1): GREEN for i in range(1, 5)}
    table(s, 0.5, 4.72, 12.33, rows, [2.9, 2.0, 7.43], sizes=(12, 12), row_h=0.42, colors=colors,
          bolds={(i, 1): True for i in range(1, 5)})
    txbox(s, 0.5, 6.85, 12.4, 0.5,
          'VERDICT: negative on the organism AND on the toxin gene - and botulism (foodborne, wound, '
          'infant-intestinal or aerosol) has no dry-surface transmission route anyway. Close the call out.',
          size=12.5, bold=True, color=GREEN, line=1.15)


def s_wbm185_evidence(prs):
    s = blank(prs)
    title(s, 'WBM185 - how mecA and CTX-M are actually called',
          'Reading the MEGARes rows in the HTML report: why MEG_3778 and MEG_2430, and not the others')
    txbox(s, 0.5, 1.38, 12.33, 0.28,
          'THE MEGARes HIERARCHY:   Type  >  Class  >  Mechanism  >  GROUP  >  MEG_ accession',
          size=13, bold=True, color=BLUE)
    txbox(s, 0.5, 1.7, 12.33, 0.5,
          'A MEG_ number is not a gene - it is ONE REFERENCE ALLELE of a gene, and MEGARes stores many '
          'near-identical alleles per gene. The rows below are NOT three mecA genes and two CTX-M genes: '
          'they are one mecA read pool and one CTX-M read pool, split across alleles.',
          size=13, line=1.18)
    rows = [['Group', 'Accession', 'Coverage', 'Depth', 'Reading'],
            ['MECA', 'MEG_3778', '90.89%', '10.51x', 'BEST ALLELE - highest breadth AND highest depth. This is the representative call.'],
            ['MECA', 'MEG_3780', '59.77%', '2.83x', 'Partial cross-mapping of the same read pool onto a homologous allele'],
            ['MECA', 'MEG_3770', '58.52%', '6.68x', 'Same - not an additional mecA gene'],
            ['MECI', 'MEG_3803', '65.50%', '5.66x', 'mecI is the REPRESSOR of the mec operon - co-occurrence supports a genuine SCCmec element (mecR1 not detected)'],
            ['CTX', 'MEG_2430', '82.88%', '7.25x', 'BEST ALLELE of the blaCTX-M family - the representative ESBL call'],
            ['CTX', 'MEG_2435', '54.14%', '2.71x', 'Partial cross-mapping onto a second CTX-M allele'],
            ['BLAZ', 'MEG_1330 / 1331', '70.65% / 64.73%', '8.17x / 4.24x', 'Staphylococcal penicillinase - ordinary carriage across the whole genus, not an S. aureus finding'],
            ['MUPA', 'MEG_4089', '90.18%', '11.70x', 'High-level mupirocin resistance (alternate IleRS) - the gene that would defeat decolonisation']]
    colors = {(1, 0): RED, (1, 4): RED, (5, 0): RED, (5, 4): RED, (8, 0): RED, (8, 4): RED}
    bolds = {(1, 0): True, (1, 4): True, (5, 0): True, (5, 4): True, (8, 0): True, (8, 4): True}
    table(s, 0.5, 2.32, 12.33, rows, [0.85, 1.7, 1.4, 1.2, 7.18], sizes=(12, 12), row_h=0.42,
          colors=colors, bolds=bolds)
    txbox(s, 0.5, 6.12, 6.05, 0.26, 'WHY 90.89% COVERAGE IS THE STRONG PART', size=12, bold=True, color=GREEN)
    txbox(s, 0.5, 6.40, 6.05, 0.95,
          'Breadth, not depth, separates a real gene from a conserved fragment. 90.89% of the ~2 kb mecA '
          'reference carries read support at 10.5x - the reason mecA is the strongest AMR call in the '
          'batch AT READ LEVEL.',
          size=12, line=1.12)
    txbox(s, 6.78, 6.12, 6.05, 0.26, 'WHY IT STILL IS NOT "MRSA"', size=12, bold=True, color=RED)
    txbox(s, 6.78, 6.40, 6.05, 1.0,
          'MRSA = methicillin-resistant Staphylococcus aureus - the SPECIES is half the term. '
          'Read-level calling carries no linkage, and S. aureus is only 0.73% of the staphylococcal '
          'reads here (2,300 of 313,026; S. hominis alone 164,436). Naming MRSA would assign mecA to '
          'the rarest candidate. Only culture with AST (antimicrobial susceptibility testing) settles it.',
          size=11.5, line=1.1)


def s_wbm232_evidence(prs):
    s = blank(prs)
    title(s, 'WBM232 - CTX-M with LpxA and AdeJ, and what XDR means',
          'Why partial coverage on the A. baumannii genome limits what can be claimed')
    rows = [['Determinant', 'Cov / Depth', 'What it is', 'What it does NOT prove'],
            ['CTX-M\nMEG_2378', '60.56%\n11.69x',
             'Acquired class A extended-spectrum beta-lactamase; destroys 3rd-generation cephalosporins. '
             'The only genuinely ACQUIRED resistance gene of the three.',
             'That it sits on the A. baumannii genome rather than on the K. pneumoniae also present '
             '(311 reads)'],
            ['AdeJ\nMEG_692', '53.72%\n5.38x',
             'Inner-membrane transporter of the AdeIJK RND efflux pump. CHROMOSOMAL AND INTRINSIC to '
             'essentially every A. baumannii; resistance comes from OVEREXPRESSION, usually via an adeN '
             'repressor mutation. The whole family is here - AdeH 88.6%/19.8x, AdeN 85.8%/8.9x, AdeG, '
             'AdeI, AdeL, AdeT1/T2.',
             'Anything about resistance - DNA cannot measure expression. It is better evidence that '
             'A. baumannii is genuinely present than that it is resistant.'],
            ['LpxA\nMEG_3626', '51.65%\n1.90x',
             'First enzyme of lipid A biosynthesis. Loss-of-function in lpxA/lpxC/lpxD abolishes LPS, and '
             'since colistin binds lipid A, an LPS-null A. baumannii is fully colistin-resistant. MEGARes '
             'files it as "colistin-resistant MUTANT".',
             'That the resistance mutation is there. lpxA is a core essential gene in every '
             'Gram-negative - presence is the default state. At 1.90x no variant can be called.']]
    colors = {(1, 0): RED, (2, 0): AMBER, (3, 0): AMBER}
    table(s, 0.5, 1.38, 12.33, rows, [1.5, 1.15, 5.85, 3.83], sizes=(12, 12), row_h=1.12,
          colors=colors, bolds={(i, 0): True for i in (1, 2, 3)})
    txbox(s, 0.5, 5.32, 6.05, 0.28, 'WHY LOW COVERAGE IS THE PROBLEM', size=13, bold=True, color=RED)
    txbox(s, 0.5, 5.64, 6.05, 1.45,
          '6,496 A. baumannii reads x 150 bp is ~0.97 Mb over a ~3.9 Mb genome - only ~0.25x '
          'genome-average. At 2-5x on one gene you cannot call the point mutations that lpxA/gyrA/parC '
          'resistance depends on, you cannot tell chromosome from plasmid from another organism because '
          'reads carry no linkage, and a sequencing error looks like a real SNP.',
          size=13, line=1.15)
    txbox(s, 6.78, 5.32, 6.05, 0.28, 'XDR - THE DEFINITION AND THE STAKES', size=13, bold=True, color=RED)
    txbox(s, 6.78, 5.64, 6.05, 1.45,
          'Magiorakos 2012 (ECDC/CDC): MDR = non-susceptible in >=3 drug categories; XDR = non-susceptible '
          "in ALL BUT <=2; PDR = all. IF one isolate had all of this - CTX-M off cephalosporins, aac(6') "
          'off aminoglycosides, AdeIJK off tetracyclines/tigecycline/fluoroquinolones, lpxA off colistin - '
          'that is XDR with the last-line drug gone. A hypothesis to culture, not a confirmed organism.',
          size=13, line=1.15)


def s_cat_a(prs):
    s = blank(prs)
    title(s, 'CDC Category A agents - are any actually present?',
          'The answer is no. Here is the evidence for each.')
    rows = [['Agent', 'Verdict', 'Evidence']]
    colors, bolds = {}, {}
    for i, (ag, vd, ev) in enumerate(CAT_A, start=1):
        rows.append([ag, vd, ev])
        colors[(i, 1)] = AMBER if vd == 'NOT ASSESSABLE' else GREEN
        bolds[(i, 1)] = True
    table(s, 0.5, 1.5, 12.33, rows, [2.31, 1.3, 8.72], sizes=(12, 12), row_h=0.72,
          colors=colors, bolds=bolds)
    txbox(s, 0.5, 6.3, 12.4, 0.7,
          'The anthrax exclusion is the important one: B. cereus group IS genuinely present in WBM174 and '
          'WBM185, above the read floor. It is excluded as B. anthracis because '
          'both virulence plasmids are entirely absent - and the one capA hit in the data belongs to the '
          'S. aureus capsule operon, not the pXO2 capsule.', size=12, color=GREY, line=1.2)


def s_cat_bc(prs):
    s = blank(prs)
    title(s, 'CDC Category B and C agents')
    rows = [['Cat', 'Agent', 'Observed', 'Assessment']]
    colors = {}
    for i, (cat, ag, obs, note) in enumerate(CAT_BC, start=1):
        rows.append([cat, ag, obs, note])
        colors[(i, 0)] = BLUE
    table(s, 0.5, 1.32, 12.33, rows, [0.4, 2.98, 2.59, 6.36], sizes=(12, 12), row_h=0.55, colors=colors)
    txbox(s, 0.5, 6.15, 12.4, 0.5,
          'Nothing in Category B or C constitutes a public-health event. The most important line is the '
          'first: the Brucella flag in all five samples is a taxonomy artifact, and taken at face value '
          'it is exactly the kind of thing that triggers an unnecessary escalation.',
          size=12, color=GREY, line=1.2)


def listed_matrix():
    """Every threat-list and watchlist organism against every sample in the batch. Derived, and
    keyed on taxid with a name fallback, so an organism that reaches the list under a renamed
    genus is not silently reported absent."""
    reps = {x: triage.load_report(x) for x in SAMPLES}
    byname, bytaxid = {}, {}
    for x, g in reps.items():
        rows = g['indentification_DNA']['speciesData']['data']
        byname[x] = {r['Scientific Name']: int(r['Real Read']) for r in rows}
        bytaxid[x] = {str(r.get('Taxid', '')).strip(): int(r['Real Read']) for r in rows}
    R = triage.RULES
    tl = {k: v for k, v in R['threat_list'].items() if isinstance(v, dict)}
    wl = {k: v for k, v in R['clinical_watchlist'].items() if isinstance(v, dict)}
    out = []
    for name, v in list(tl.items()) + list(wl.items()):
        tid = str(v.get('taxid', '')).strip()
        counts = {x: (byname[x].get(name) or bytaxid[x].get(tid) or 0) for x in SAMPLES}
        grp = f"CDC {v['tier']}" if name in tl else (v.get('priority') or 'WHO')
        out.append({'name': name, 'grp': grp, 'genome': v.get('genome', 'DNA'),
                    'counts': counts, 'total': sum(counts.values()),
                    'testable': v.get('genome', 'DNA') != 'RNA'})
    return out


def s_listed_matrix(prs):
    s = blank(prs)
    rows_all = listed_matrix()
    present = [r for r in rows_all if r['total']]
    neg = [r for r in rows_all if not r['total'] and r['testable']]
    nt = [r for r in rows_all if not r['total'] and not r['testable']]
    title(s, f'All {len(rows_all)} listed pathogens, screened across the batch',
          'CDC Category A/B/C threat list (46) + WHO priority / ESKAPE watchlist (24). '
          'Real Read counts; 0 shown as blank.')
    txbox(s, 0.5, 1.28, 12.33, 0.24,
          f'{len(present)} PRESENT   |   {len(neg)} genuinely ABSENT (testable)   |   '
          f'{len(nt)} NOT TESTED (RNA genomes, no RNA library)',
          size=12, bold=True, color=BLUE)

    ORD = {'CDC A': 0, 'CDC B': 1, 'CDC C': 2}
    def key(r):
        return (ORD.get(r['grp'], 3 + ('crit' not in r['grp']) + ('high' not in r['grp'])),
                -r['total'])
    present.sort(key=key)
    cdc = [r for r in present if r['grp'].startswith('CDC')]
    who = [r for r in present if not r['grp'].startswith('CDC')]
    allp = cdc + who
    cut = (len(allp) + 1) // 2
    half, rest = allp[:cut], allp[cut:]

    SHORTG = {'WHO critical': 'WHO crit', 'WHO critical (fungal)': 'WHO crit',
              'WHO high': 'WHO high', 'WHO medium': 'WHO med',
              'WHO (fungal)': 'WHO', 'WHO (mycobacterial)': 'WHO'}
    def draw(x0, w, data, hdr):
        txbox(s, x0, 1.62, w, 0.24, hdr, size=12, bold=True, color=RED)
        rows = [['Organism', 'List'] + [x[3:] for x in SAMPLES]]
        colors = {}
        for i, r in enumerate(data, start=1):
            colors[(i, 1)] = BLUE
            if r['grp'].startswith('CDC'):
                colors[(i, 0)] = RED
            rows.append([abbr(r['name'])[:26], SHORTG.get(r['grp'], r['grp'])]
                        + [f"{r['counts'][x]:,}" if r['counts'][x] else '' for x in SAMPLES])
        table(s, x0, 1.92, w, rows, [w - 3.7, 0.95, 0.55, 0.55, 0.55, 0.55, 0.55],
              sizes=(12, 12), row_h=0.28, colors=colors)
    draw(0.5, 5.95, half, 'PRESENT — CDC agents, then WHO critical')
    draw(6.88, 5.95, rest, 'PRESENT — WHO high / medium / other')

    y = 1.92 + 0.28 * (max(len(half), len(rest)) + 1) + 0.18
    txbox(s, 0.5, y, 5.95, 0.24, f'GENUINELY ABSENT — {len(neg)}, and the assay could have seen them',
          size=12, bold=True, color=GREEN)
    txbox(s, 0.5, y + 0.28, 5.95, 1.2, ', '.join(abbr(r['name']) for r in neg) + '.',
          size=12, color=INK, line=1.08)

    txbox(s, 6.88, y, 5.95, 0.24, f'NOT TESTED — {len(nt)}. THIS IS NOT A NEGATIVE.',
          size=12, bold=True, color=AMBER)
    txbox(s, 6.88, y + 0.28, 5.95, 1.2,
          'Every VHF, the encephalitides, influenza, SARS/MERS, Nipah, Hendra, yellow fever, '
          'chikungunya and TBE have RNA genomes. No RNA library here, so the test did not run and '
          'absence carries no information — the largest coverage gap in the assay.',
          size=12, color=INK, line=1.08)


def s_evidence_strength(prs):
    s = blank(prs)
    title(s, 'How strong is each negative?',
          'All 46 CDC agents are screened - but "not found" does not mean the same thing four times over')
    rows = [['Strength', 'n', 'What was actually done', 'Agents'],
            ['MARKER-CONFIRMED\nNEGATIVE', '9',
             'A toxin or virulence gene unique to the agent was searched for and not found. The gene is '
             'reliably in VFDB, so the miss IS the negative. These downgrade the organism outright.',
             'anthrax (pXO1/pXO2), botulism (bont), plague (caf1/lcrV/pla), tularaemia (fopA/tul4), '
             'C. perfringens (etx), S. aureus (seb), cholera (ctxA/ctxB), Shigella and E. coli (stx1/stx2/eae)'],
            ['MARKER-SUPPORTED\nNEGATIVE', '8',
             'A marker was searched for and not found, but VFDB coverage for the agent cannot be verified '
             'from the report. Finding it would escalate; NOT finding it changes nothing. A one-way test.',
             'Brucella x3 (btpA/btpB, omp25/31, bvrR/S), Burkholderia x2 (bsa T3SS, bimA, capsule), '
             'Coxiella burnetii (dotA/dotB), Salmonella (Vi capsule tviA-E), M. tuberculosis (esxA/esxB)'],
            ['TAXONOMY ONLY', '4',
             'No usable marker exists. Rests on the read assignment and on naming the look-alikes it '
             'could be confused with.',
             'Variola (VFDB is bacterial), C. psittaci (its genes are genus-wide), R. prowazekii (the '
             'discriminator is an ABSENCE - typhus group lacks ompA), Cryptosporidium (not a bacterium)'],
            ['NOT ASSESSABLE', '25',
             'RNA genomes against a DNA-only library. The test did not run. This is a missing test, not '
             'a negative result, and it is never allowed to collapse into one.',
             'All viral haemorrhagic fevers, the encephalitides, influenza, SARS/MERS, Nipah, Hendra, '
             'yellow fever, chikungunya, TBE']]
    colors = {(1, 0): GREEN, (2, 0): BLUE, (3, 0): AMBER, (4, 0): RED}
    bolds = {(i, 0): True for i in (1, 2, 3, 4)}
    table(s, 0.5, 1.45, 12.33, rows, [1.62, 0.42, 4.44, 5.85], sizes=(12, 11), row_h=1.05,
          colors=colors, bolds=bolds)
    txbox(s, 0.5, 6.18, 12.4, 0.9,
          'WHY THE MIDDLE BAND IS ONE-WAY: a marker that is not in the reference database can never be '
          'found, and a two-way test would read that as a clean bill of health. Treating the miss as no '
          'information is the only safe direction.\n'
          'Marker searches are also restricted to reference strains of the same genus - unrestricted, '
          'S. aureus\'s own esxA/esxB would have "confirmed" M. tuberculosis on a certified-clean control.',
          size=12, color=GREY, line=1.18, space_after=4)


def s_amr(prs, AM):
    s = blank(prs)
    title(s, 'Antimicrobial resistance - where the real signal is',
          'MEGARes gene hits; coverage is breadth of the reference gene, depth is mean per-base coverage')
    KEY = {
        'WBM156': '- none of concern (13 classes, all intrinsic ribosomal / efflux)',
        'WBM174': 'blaZ 11.71x, ermC 10.17x  - community-normal S. aureus carriage',
        'WBM179': 'blaZ 6.04x, mecI 1.62x  - mecI regulator WITHOUT mecA; not MRSA',
        'WBM185': "mecA 90.9% / 10.51x,  mupA 90.2% / 11.70x,  CTX-M 82.9% / 7.25x,  aac(6') 79.8% / 7.42x,  blaZ 70.7% / 8.17x",
        'WBM232': "CTX-M 60.6% / 11.69x,  aac(6') 63.7% / 14.08x,  AdeH 88.6% / 19.8x,  AdeN 85.8% / 8.85x,  AdeJ 53.7% / 5.38x,  LpxA 51.7% / 1.90x",
    }
    rows = [['Sample', 'AMR genes', 'Classes', 'Key beta-lactam / high-risk determinants']]
    colors, bolds = {}, {}
    for i, sm in enumerate(SAMPLES, start=1):
        rows.append([sm, AM[sm][0], AM[sm][1], KEY[sm]])
        colors[(i, 0)] = BLUE
        if sm in ('WBM185', 'WBM232'):
            colors[(i, 3)] = RED; bolds[(i, 3)] = True
    table(s, 0.5, 1.5, 12.33, rows, [1.1, 1.1, 0.9, 9.22], sizes=(12, 12), row_h=0.42,
          colors=colors, bolds=bolds)
    txbox(s, 0.5, 4.22, 12.4, 0.28, 'THE LIMITATION THAT MATTERS', size=13, bold=True, color=RED)
    txbox(s, 0.5, 4.57, 12.4, 2.3,
          'The PFI pipeline calls AMR genes straight off reads, so a gene is never linked to the organism '
          'carrying it. "mecA present" in a sample containing both S. aureus (2,300 reads) AND abundant '
          'coagulase-negative staphylococci (S. hominis 11.2%, S. haemolyticus 1.8%, S. epidermidis 4.1%) '
          'does not tell you whose mecA it is - and methicillin-resistant S. epidermidis is an ordinary '
          'skin organism, whereas MRSA on a check-in kiosk is a different conversation.\n\n'
          'Nothing in this report can close that gap. The resistance table has no organism column, and '
          'the species table has no gene column - there is no field in either that joins them. The host '
          'of mecA, CTX-M, lpxA and adeJ is genuinely unresolved, and culture with AST is the way to '
          'settle it.',
          size=13, line=1.22, space_after=8)


def s_recs(prs):
    s = blank(prs)
    title(s, 'Recommendations')
    # Three items, full width, so point 1 has room to show its arithmetic. Points on the MRSA
    # wording and on reporting Real Read were removed 2026-08-12 at the user's request; the MRSA
    # caveat is carried in full on the WBM185 sample and evidence slides, and the Real Read point
    # is the first of the three traps on the method slide, so neither claim is lost.
    recs = [
        ('1', 'Re-swab WBM232 (T4 trolley handles) with a culture arm', RED,
         'Metagenomics cannot produce an antibiogram. It also cannot call the point mutations that '
         'matter here, and the arithmetic shows why: 6,496 A. baumannii reads x 150 bp = ~0.97 Mb '
         'spread over a ~3.9 Mb genome, so the organism is sequenced at only ~0.25x on average - '
         'roughly one base in four is seen ONCE, and most positions are not seen at all.\n'
         'gyrA (99.5% cov, 120.7x) and lpxA (51.7% cov, 1.9x) ARE both detected in WBM232. That '
         'changes nothing: every A. baumannii carries them, so PRESENCE is the default state and '
         'resistance needs a specific MUTATION. The gyrA depth is community-wide - gyrA is a '
         'housekeeping gene in essentially every bacterium in the swab, and A. baumannii contributes '
         'a small fraction of those reads, which cannot be separated out. An isolate is needed.'),
        ('2', 'Add an RNA library to the protocol', AMBER,
         'Four of six Category A agents and every respiratory virus of surveillance interest are RNA '
         'viruses. The current assay cannot see them at all - they are reported NOT_TESTED, never '
         'negative. This is the single largest gap in the batch.'),
        ('3', 'Include a blank extraction control in every batch', AMBER,
         'Without one, no trace call below ~100 reads can be formally cleared as reagent background. '
         'A blank would have settled the botulism and cholera calls in minutes rather than by '
         'cross-sample argument.'),
    ]
    y = 1.5
    for n, head, rgb, body in recs:
        c = s.shapes.add_shape(9, Inches(0.5), Inches(y), Inches(0.36), Inches(0.36))
        c.fill.solid(); c.fill.fore_color.rgb = rgb
        c.line.fill.background(); c.shadow.inherit = False
        pgh = c.text_frame.paragraphs[0]; pgh.alignment = PP_ALIGN.CENTER
        r = pgh.add_run(); r.text = n
        _style(r, 13, True, WHITE)
        txbox(s, 1.0, y - 0.02, 11.83, 0.3, head, size=13.5, bold=True, color=rgb)
        # Advance by what this item actually occupies, not a fixed step: a fixed step left a
        # half-slide gap under the short items and read as a missing recommendation.
        lines = sum(max(1, -(-len(para) // 150)) for para in body.split('\n'))
        h = lines * 0.20
        txbox(s, 1.0, y + 0.3, 11.83, h + 0.1, body, size=12, line=1.08, space_after=2)
        y += 0.34 + h + 0.34


def build():
    prs = Presentation(TPL)
    Q, SP, AM, VF = qc(), species(), amr_counts(), vf_counts()

    s_cover(prs)
    s_method(prs)
    s_qc(prs, Q)
    s_caveats(prs)
    s_topspecies(prs, SP)
    s_howtoread(prs)
    for sm in SAMPLES:
        s_sample(prs, sm, Q, SP, AM, VF)
        if sm == 'WBM174':
            s_botulism(prs)
        elif sm == 'WBM185':
            s_wbm185_evidence(prs)
        elif sm == 'WBM232':
            s_wbm232_evidence(prs)
    s_cat_a(prs)
    s_cat_bc(prs)
    s_listed_matrix(prs)
    s_evidence_strength(prs)
    s_amr(prs, AM)
    s_recs(prs)

    # template slides: 0=cover, 1=empty placeholder, 2="Thank you"; ours start at 3.
    lst = prs.slides._sldIdLst
    ids = list(lst)
    lst.remove(ids[1])
    lst.remove(ids[2])
    lst.append(ids[2])

    prs.save(OUT)
    print(f'{OUT}: {len(list(lst))} slides')


if __name__ == '__main__':
    build()
